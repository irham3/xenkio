"""
Document-to-Markdown converter for Pyodide.

Provides two conversion strategies:
  - convert_lazy(file_path, file_type)  — uses lightweight per-format libs
  - init_preload() / convert_preload(file_path) — uses markitdown (all formats)
"""

import csv
import json
import sys


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _rows_to_markdown_table(rows: list[list[str]]) -> str:
    """Convert a list of rows into a Markdown table string."""
    if not rows:
        return ""

    lines = []
    header = rows[0]
    lines.append("| " + " | ".join(str(c) if c is not None else "" for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
    return "\n".join(lines)


def _read_text(file_path: str) -> str:
    """Read a file as UTF-8 text, ignoring decode errors."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ---------------------------------------------------------------------------
# Per-format converters (lazy strategy)
# ---------------------------------------------------------------------------

def _convert_pdf(file_path: str) -> str:
    from pdfminer.high_level import extract_text  # type: ignore
    return extract_text(file_path)


def _convert_docx(file_path: str) -> str:
    from docx import Document  # type: ignore
    doc = Document(file_path)
    md_lines = []
    
    def extract_run_text(run):
        text = run.text
        if not text.strip():
            return text
        l_space = len(text) - len(text.lstrip())
        r_space = len(text) - len(text.rstrip())
        clean_text = text.strip()
        if run.bold:
            clean_text = f"**{clean_text}**"
        if run.italic:
            clean_text = f"*{clean_text}*"
        if run.underline:
            clean_text = f"<u>{clean_text}</u>"
        return (" " * l_space) + clean_text + (" " * r_space)
        
    for element in doc.element.body:
        if element.tag.endswith('p'):
            from docx.text.paragraph import Paragraph  # type: ignore
            p = Paragraph(element, doc)
            text = ""
            for run in p.runs:
                text += extract_run_text(run)
            
            style_name = p.style.name.lower()
            if 'heading' in style_name:
                try:
                    level = int(''.join(filter(str.isdigit, style_name)))
                    if level > 0 and level <= 6:
                        text = f"{'#' * level} {text}"
                except ValueError:
                    pass
            elif style_name == 'title':
                text = f"# {text}"
            elif style_name == 'subtitle':
                text = f"## {text}"
                
            if 'list paragraph' in style_name or p.style.name == 'List Paragraph':
                text = f"- {text}"
                
            if text.strip() or len(text) > 0:
                md_lines.append(text)
                
        elif element.tag.endswith('tbl'):
            from docx.table import Table  # type: ignore
            t = Table(element, doc)
            table_lines = []
            for r_idx, row in enumerate(t.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = ""
                    for cell_p in cell.paragraphs:
                        p_text = ""
                        for run in cell_p.runs:
                            p_text += extract_run_text(run)
                        cell_text += p_text + "<br>"
                    if cell_text.endswith("<br>"):
                        cell_text = cell_text[:-4]
                    cell_text = cell_text.replace("\n", " ").strip()
                    row_data.append(cell_text)
                
                table_lines.append("| " + " | ".join(row_data) + " |")
                if r_idx == 0:
                    table_lines.append("| " + " | ".join(["---"] * len(row_data)) + " |")
            md_lines.append("\n".join(table_lines) + "\n")
            
    return "\n\n".join(md_lines)


def _convert_xlsx(file_path: str) -> str:
    import openpyxl  # type: ignore
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        sections.append(f"# {sheet_name}\n")
        sections.append(_rows_to_markdown_table(rows))
        sections.append("")
    return "\n".join(sections)


def _convert_pptx(file_path: str) -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation(file_path)
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        sections.append(f"## Slide {i}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                sections.append(shape.text)
        sections.append("")
    return "\n".join(sections)


def _convert_html(file_path: str) -> str:
    from bs4 import BeautifulSoup  # type: ignore
    raw = _read_text(file_path)
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text("\n", strip=True)


def _convert_csv(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        rows = list(csv.reader(f))
    return _rows_to_markdown_table(rows)


def _convert_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return "```json\n" + json.dumps(data, indent=2) + "\n```"


def _convert_rtf(file_path: str) -> str:
    from striprtf.striprtf import rtf_to_text  # type: ignore
    return rtf_to_text(_read_text(file_path))


# Extension → converter mapping
_LAZY_CONVERTERS: dict[str, callable] = {
    "pdf":  _convert_pdf,
    "docx": _convert_docx,
    "xlsx": _convert_xlsx,
    "xls":  _convert_xlsx,
    "pptx": _convert_pptx,
    "html": _convert_html,
    "htm":  _convert_html,
    "csv":  _convert_csv,
    "json": _convert_json,
    "rtf":  _convert_rtf,
}


def convert_lazy(file_path: str, file_type: str) -> str:
    """Convert a document using format-specific libraries (lazy strategy)."""
    converter = _LAZY_CONVERTERS.get(file_type)
    if converter:
        return converter(file_path)
    # Fallback: read as plain text (txt, md, rst, log, xml, etc.)
    return _read_text(file_path)


# ---------------------------------------------------------------------------
# markitdown-based converter (preload strategy)
# ---------------------------------------------------------------------------

_markitdown_instance = None


def init_preload() -> None:
    """
    Initialize markitdown after all format libraries have been installed.

    Must be called *after* micropip has installed all dependencies, because
    markitdown's converter modules check for them at import time and cache
    the result in `_dependency_exc_info` module-level variables.
    """
    global _markitdown_instance

    # 1. Purge all cached markitdown modules so import-time checks re-run
    stale_keys = [k for k in sys.modules if "markitdown" in k]
    for k in stale_keys:
        del sys.modules[k]

    # 2. Fresh import
    from markitdown import MarkItDown
    import markitdown

    # 3. Safety net: clear any stale _dependency_exc_info in converter modules
    for mod_name, mod_obj in list(sys.modules.items()):
        if "markitdown" in mod_name and hasattr(mod_obj, "_dependency_exc_info"):
            mod_obj._dependency_exc_info = None

    _markitdown_instance = MarkItDown()


def convert_preload(file_path: str) -> str:
    """Convert a document using markitdown (preload strategy), with fallback to lazy."""
    if _markitdown_instance is None:
        raise RuntimeError("init_preload() must be called before convert_preload()")
    
    # Intercept docx to use our custom robust converter
    ext = file_path.split('.')[-1].lower()
    if ext == 'docx':
        return _convert_docx(file_path)
    
    try:
        result = _markitdown_instance.convert(file_path)
        return result.text_content
    except Exception as e:
        # MarkItDown sometimes fails in Pyodide due to complex dependencies like pandas.
        # If it fails, fallback to our lightweight, pure-Python lazy converters.
        if ext in _LAZY_CONVERTERS:
            print(f"MarkItDown failed for .{ext}, falling back to lazy converter. Error: {e}")
            return convert_lazy(file_path, ext)
        raise e
